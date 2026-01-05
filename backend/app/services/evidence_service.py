"""Evidence service for document processing and analysis."""
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any
from uuid import UUID

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.models.evidence import Evidence
from app.models.assessment import AssessmentResponse
from app.models.ndi import NDIQuestion, NDIMaturityLevel
from app.models.settings import SystemPrompt
from app.schemas.evidence import EvidenceAnalysis
from app.config import settings
from app.services.ai_service import get_active_ai_provider


async def get_system_prompt(db, prompt_id: str) -> Optional[str]:
    """Get a system prompt template from the database."""
    result = await db.execute(
        select(SystemPrompt).where(SystemPrompt.id == prompt_id).where(SystemPrompt.is_active == True)
    )
    prompt = result.scalar_one_or_none()
    return prompt.prompt_template if prompt else None


class EvidenceService:
    """Service for evidence processing and analysis."""

    def __init__(self, db: AsyncSession):
        self.db = db

    async def extract_text(self, evidence: Evidence) -> Optional[str]:
        """Extract text from uploaded document."""
        file_path = Path(evidence.file_path)

        if not file_path.exists():
            return None

        extracted_text = ""
        file_type = evidence.file_type.lower() if evidence.file_type else ""

        try:
            if file_type == "pdf":
                extracted_text = await self._extract_pdf(file_path)
            elif file_type in ["docx", "doc"]:
                extracted_text = await self._extract_docx(file_path)
            elif file_type in ["xlsx", "xls"]:
                extracted_text = await self._extract_excel(file_path)
            elif file_type in ["pptx", "ppt"]:
                extracted_text = await self._extract_pptx(file_path)
            elif file_type == "txt":
                extracted_text = file_path.read_text(encoding="utf-8")

            # Update evidence record
            evidence.extracted_text = extracted_text
            await self.db.flush()

            return extracted_text
        except Exception as e:
            print(f"Error extracting text: {e}")
            return None

    async def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
        except Exception as e:
            print(f"PDF extraction error: {e}")
            return ""

    async def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX."""
        try:
            from docx import Document
            doc = Document(str(file_path))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            return text.strip()
        except Exception as e:
            print(f"DOCX extraction error: {e}")
            return ""

    async def _extract_excel(self, file_path: Path) -> str:
        """Extract text from Excel."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows():
                    row_text = " ".join(
                        str(cell.value) if cell.value else "" for cell in row
                    )
                    if row_text.strip():
                        text += row_text + "\n"
            return text.strip()
        except Exception as e:
            print(f"Excel extraction error: {e}")
            return ""

    async def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint."""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            print(f"PPTX extraction error: {e}")
            return ""

    async def analyze_evidence(self, evidence_id: UUID) -> EvidenceAnalysis:
        """Analyze evidence using AI."""
        # Get evidence with response and question
        result = await self.db.execute(
            select(Evidence)
            .options(
                selectinload(Evidence.response)
                .selectinload(AssessmentResponse.question)
                .selectinload(NDIQuestion.maturity_levels)
            )
            .where(Evidence.id == evidence_id)
        )
        evidence = result.scalar_one_or_none()

        if not evidence:
            raise ValueError("Evidence not found")

        # Get the selected level details
        selected_level = evidence.response.selected_level if evidence.response else 0
        question = evidence.response.question if evidence.response else None

        if not question:
            return EvidenceAnalysis(
                supports_level="no",
                covered_criteria=[],
                missing_criteria=[],
                confidence_score=0.0,
                recommendations=["Cannot analyze without question context"],
            )

        # Get maturity level criteria
        level_criteria = next(
            (ml for ml in question.maturity_levels if ml.level == selected_level),
            None
        )

        if not level_criteria:
            return EvidenceAnalysis(
                supports_level="no",
                covered_criteria=[],
                missing_criteria=[],
                confidence_score=0.0,
                recommendations=["Maturity level criteria not found"],
            )

        # Perform AI analysis
        analysis = await self._ai_analyze(
            document_text=evidence.extracted_text or "",
            question=question.question_ar,
            level_description=level_criteria.description_ar,
            acceptance_criteria=level_criteria.acceptance_evidence_ar or [],
        )

        # Update evidence record
        evidence.ai_analysis = analysis
        evidence.analysis_status = "completed"
        evidence.analyzed_at = datetime.utcnow()
        await self.db.flush()

        return EvidenceAnalysis(
            supports_level=analysis.get("supports_level", "no"),
            covered_criteria=analysis.get("covered_criteria", []),
            missing_criteria=analysis.get("missing_criteria", []),
            confidence_score=analysis.get("confidence_score", 0.0),
            recommendations=analysis.get("recommendations", []),
            summary_ar=analysis.get("summary_ar"),
            summary_en=analysis.get("summary_en"),
        )

    async def analyze_evidence_against_criteria(
        self,
        evidence_id: UUID,
        question_code: str,
        selected_level: int,
        language: str = "ar",
    ) -> EvidenceAnalysis:
        """Analyze evidence against specific criteria."""
        # Get evidence
        evidence_result = await self.db.execute(
            select(Evidence).where(Evidence.id == evidence_id)
        )
        evidence = evidence_result.scalar_one_or_none()

        if not evidence:
            raise ValueError("Evidence not found")

        # Get question with levels
        question_result = await self.db.execute(
            select(NDIQuestion)
            .options(selectinload(NDIQuestion.maturity_levels))
            .where(NDIQuestion.code == question_code.upper())
        )
        question = question_result.scalar_one_or_none()

        if not question:
            raise ValueError("Question not found")

        # Get level criteria
        level_criteria = next(
            (ml for ml in question.maturity_levels if ml.level == selected_level),
            None
        )

        if not level_criteria:
            return EvidenceAnalysis(
                supports_level="no",
                covered_criteria=[],
                missing_criteria=[],
                confidence_score=0.0,
                recommendations=["Maturity level criteria not found"],
            )

        # Use appropriate language
        if language == "ar":
            question_text = question.question_ar
            level_desc = level_criteria.description_ar
            criteria = level_criteria.acceptance_evidence_ar or []
        else:
            question_text = question.question_en
            level_desc = level_criteria.description_en
            criteria = level_criteria.acceptance_evidence_en or []

        # Perform analysis
        analysis = await self._ai_analyze(
            document_text=evidence.extracted_text or "",
            question=question_text,
            level_description=level_desc,
            acceptance_criteria=criteria,
        )

        # Update evidence
        evidence.ai_analysis = analysis
        evidence.analysis_status = "completed"
        evidence.analyzed_at = datetime.utcnow()
        await self.db.flush()

        return EvidenceAnalysis(
            supports_level=analysis.get("supports_level", "no"),
            covered_criteria=analysis.get("covered_criteria", []),
            missing_criteria=analysis.get("missing_criteria", []),
            confidence_score=analysis.get("confidence_score", 0.0),
            recommendations=analysis.get("recommendations", []),
            summary_ar=analysis.get("summary_ar"),
            summary_en=analysis.get("summary_en"),
        )

    async def _ai_analyze(
        self,
        document_text: str,
        question: str,
        level_description: str,
        acceptance_criteria: list[str],
    ) -> dict:
        """Perform AI analysis of document using database-configured provider."""
        # Get AI provider from database
        ai_provider = await get_active_ai_provider(self.db)

        if not ai_provider or not ai_provider.get("api_key"):
            # Return mock analysis if no AI configured
            return {
                "supports_level": "partial",
                "covered_criteria": acceptance_criteria[:len(acceptance_criteria)//2] if acceptance_criteria else [],
                "missing_criteria": acceptance_criteria[len(acceptance_criteria)//2:] if acceptance_criteria else [],
                "confidence_score": 0.5,
                "recommendations": [
                    "Configure AI providers in Settings for detailed analysis",
                    "Review document manually against criteria",
                ],
                "summary_ar": "تحليل أولي - يرجى تكوين مزود الذكاء الاصطناعي في الإعدادات للتحليل التفصيلي",
                "summary_en": "Preliminary analysis - Please configure AI provider in Settings for detailed analysis",
            }

        try:
            return await self._analyze_with_provider(
                ai_provider, document_text, question, level_description, acceptance_criteria
            )
        except Exception as e:
            print(f"AI analysis error: {e}")
            return {
                "supports_level": "no",
                "covered_criteria": [],
                "missing_criteria": acceptance_criteria,
                "confidence_score": 0.0,
                "recommendations": [f"Analysis failed: {str(e)}"],
            }

    async def _analyze_with_provider(
        self,
        ai_provider: Dict[str, Any],
        document_text: str,
        question: str,
        level_description: str,
        acceptance_criteria: list[str],
    ) -> dict:
        """Analyze using the configured AI provider."""
        import json

        provider_id = ai_provider.get("id", "")
        api_key = ai_provider.get("api_key", "")
        model_name = ai_provider.get("model_name", "")
        api_endpoint = ai_provider.get("api_endpoint", "")

        criteria_text = "\n".join(f"- {c}" for c in acceptance_criteria)

        # Try to get prompt from database, fallback to default
        prompt_template = await get_system_prompt(self.db, "evidence_analysis")
        if prompt_template:
            prompt = prompt_template.format(
                question=question,
                level_description=level_description,
                criteria_text=criteria_text,
                document_text=document_text[:8000]
            )
        else:
            # Fallback default prompt
            prompt = f"""أنت محلل شواهد لمؤشر البيانات الوطني (NDI).

السؤال: {question}
وصف المستوى: {level_description}
معايير القبول:
{criteria_text}

المستند المرفوع:
{document_text[:8000]}

المطلوب:
1. هل المستند يدعم المستوى المختار؟ (yes/partial/no)
2. أي معايير قبول يغطيها المستند؟
3. أي معايير قبول مفقودة؟
4. توصيات لتحسين الشواهد

أجب بـ JSON فقط بدون أي نص إضافي:
{{
  "supports_level": "yes|partial|no",
  "covered_criteria": [],
  "missing_criteria": [],
  "confidence_score": 0.0-1.0,
  "recommendations": [],
  "summary_ar": "",
  "summary_en": ""
}}
"""

        response_text = ""

        if provider_id == "gemini":
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            model = genai.GenerativeModel(model_name or "gemini-pro")
            response = model.generate_content(prompt)
            response_text = response.text

        elif provider_id == "openai":
            from openai import AsyncOpenAI
            client = AsyncOpenAI(api_key=api_key)
            response = await client.chat.completions.create(
                model=model_name or "gpt-4",
                messages=[{"role": "user", "content": prompt}],
                response_format={"type": "json_object"},
            )
            response_text = response.choices[0].message.content

        elif provider_id == "claude":
            import anthropic
            client = anthropic.AsyncAnthropic(api_key=api_key)
            response = await client.messages.create(
                model=model_name or "claude-3-opus-20240229",
                max_tokens=2048,
                messages=[{"role": "user", "content": prompt}],
            )
            response_text = response.content[0].text

        elif provider_id == "azure":
            from openai import AsyncAzureOpenAI
            client = AsyncAzureOpenAI(
                api_key=api_key,
                api_version="2024-02-15-preview",
                azure_endpoint=api_endpoint,
            )
            response = await client.chat.completions.create(
                model=model_name or "gpt-4",
                messages=[{"role": "user", "content": prompt}],
            )
            response_text = response.choices[0].message.content
        else:
            raise ValueError(f"Unknown AI provider: {provider_id}")

        # Parse JSON response
        try:
            # Find JSON in response
            start = response_text.find("{")
            end = response_text.rfind("}") + 1
            if start >= 0 and end > start:
                return json.loads(response_text[start:end])
        except json.JSONDecodeError:
            pass

        return {
            "supports_level": "partial",
            "covered_criteria": [],
            "missing_criteria": acceptance_criteria,
            "confidence_score": 0.5,
            "recommendations": ["Could not parse AI response"],
        }
